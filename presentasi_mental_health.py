from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
import os

# --- DIREKTORI UNTUK PLOT ---
# Pastikan folder ini ada dan berisi gambar-gambar Anda
plots_dir = "plots_output"

# --- DEFENISI WARNA DAN FONT UNTUK PPT ---
TEXT_COLOR = RGBColor(60, 60, 60) # Abu-abu gelap untuk teks utama
ACCENT_COLOR = RGBColor(0, 102, 204) # Biru aksen (IBM Blue)
LIGHT_BACKGROUND_COLOR = RGBColor(240, 248, 255) # Latar belakang terang (AliceBlue)
FOOTER_COLOR = RGBColor(120, 120, 120) # Abu-abu untuk teks footer

BODY_FONT = "Segoe UI"
TITLE_FONT = "Segoe UI Light"

# --- FUNGSI-FUNGSI BANTUAN UNTUK PPT ---

def set_text_properties(text_frame, content_text, font_size, bold=False, color=TEXT_COLOR, font_name=BODY_FONT, align_center=False, is_bullet=False):
    """
    Mengatur properti teks untuk text_frame yang diberikan.
    Mendukung multiline string, perataan, dan bullet points.
    """
    text_frame.clear()
    paragraphs = content_text.split('\n')

    for i, line in enumerate(paragraphs):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        run = p.add_run()
        
        if is_bullet and line.strip().startswith('- '):
            p.level = 0
            run.text = line.strip()[2:] # Hapus tanda '-'
        else:
            run.text = line.strip()

        font = run.font
        font.size = font_size
        font.bold = bold
        font.name = font_name
        font.color.rgb = color

        if align_center:
            p.alignment = PP_ALIGN.CENTER
        else:
            p.alignment = PP_ALIGN.LEFT

def add_simple_background(slide, color=LIGHT_BACKGROUND_COLOR):
    """Menambahkan bentuk persegi panjang sebagai latar belakang slide."""
    left = top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    line = shape.line
    line.fill.background()

def add_footer(slide, text_content="IBM Granite - Mental Health in Tech", slide_number=None):
    """Menambahkan footer di bagian bawah slide."""
    left = Inches(0.5)
    top = Inches(prs.slide_height - Inches(0.5))
    width = prs.slide_width - Inches(1)
    height = Inches(0.3)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    
    p_left = text_frame.paragraphs[0]
    run_left = p_left.add_run()
    run_left.text = text_content
    run_left.font.size = Pt(10)
    run_left.font.name = BODY_FONT
    run_left.font.color.rgb = FOOTER_COLOR
    p_left.alignment = PP_ALIGN.LEFT

    if slide_number is not None:
        p_right = text_frame.add_paragraph()
        run_right = p_right.add_run()
        run_right.text = f"Slide {slide_number}"
        run_right.font.size = Pt(10)
        run_right.font.name = BODY_FONT
        run_right.font.color.rgb = FOOTER_COLOR
        p_right.alignment = PP_ALIGN.RIGHT
        p_right.space_before = Pt(0)
        p_right.space_after = Pt(0)


# --- INISIALISASI PRESENTASI ---
print("--- Memulai Pembuatan Presentasi PPT ---")
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

# --- Variabel Hasil Analisis (ANDA BISA MENGGANTI INI DENGAN HASIL AKTUAL ANDA) ---
# Gunakan nilai ROC AUC yang Anda dapatkan dari analisis sebelumnya
roc_score_value = 0.88 # Contoh nilai, ganti dengan nilai aktual Anda

# Gunakan daftar fitur paling berpengaruh yang Anda dapatkan
top_features_list_value = ["work_interfere", "family_history", "anonymity"]

# GANTI DENGAN RINGKASAN EKSEKUTIF AKTUAL DARI REPLICATE API ANDA
# Jika Anda tidak menggunakan Replicate, biarkan teks ini atau sesuaikan.
final_executive_summary = """
Model Random Forest berhasil memprediksi apakah pekerja pernah menerima
perawatan mental berdasarkan data survei. Fitur paling berpengaruh adalah
interferensi kerja, riwayat keluarga, dan anonimitas.

Sebanyak 51% responden pernah mendapat perawatan mental. Skor ROC AUC sebesar
0.88 menunjukkan keandalan prediksi yang tinggi, menegaskan model ini dapat diandalkan.

Insight ini krusial bagi manajemen HR untuk mengembangkan kebijakan dan program
dukungan kesehatan mental yang lebih bertarget, meningkatkan kesejahteraan
karyawan, dan mengurangi dampak stres kerja di sektor teknologi.
"""

# --- PATH GAMBAR (PASTIKAN GAMBAR ADA DI FOLDER 'plots_output') ---
treatment_plot_path = os.path.join(plots_dir, "treatment_distribution.png")
confusion_matrix_path = os.path.join(plots_dir, "confusion_matrix.png")
feature_importance_path = os.path.join(plots_dir, "feature_importance.png")

# --- PEMBUATAN SLIDE-SLIDE PPT ---

# Slide 1: Judul
slide = prs.slides.add_slide(title_slide_layout)
add_simple_background(slide, color=LIGHT_BACKGROUND_COLOR)

title_shape = slide.shapes.title
subtitle_shape = slide.placeholders[1]

set_text_properties(title_shape.text_frame, "🧠 Mental Health in Tech", Pt(54), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT, align_center=True)
title_shape.top = Inches(2)
title_shape.left = Inches(1)
title_shape.width = Inches(8)
title_shape.height = Inches(2)

set_text_properties(subtitle_shape.text_frame, "Prediction & Summarization\nCapstone Project – IBM Granite via Replicate", Pt(30), bold=False, color=TEXT_COLOR, font_name=BODY_FONT, align_center=True)
subtitle_shape.top = Inches(4.5)
subtitle_shape.left = Inches(1)
subtitle_shape.width = Inches(8)
subtitle_shape.height = Inches(1.5)


# Slide 2: Dataset Overview
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "1. Dataset Overview", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

content_placeholder = slide.placeholders[1]
content_text = (
    "- Dataset: Mental Health in Tech Survey (OSMI)\n"
    "- Jumlah fitur: ~25\n"
    "- Target klasifikasi: Apakah seseorang pernah menerima perawatan mental\n"
    "- Data dibersihkan & dikodekan secara kategorikal"
)
set_text_properties(content_placeholder.text_frame, content_text, Pt(24), color=TEXT_COLOR, font_name=BODY_FONT, is_bullet=True)
content_placeholder.top = Inches(2)
content_placeholder.left = Inches(1)
content_placeholder.width = Inches(8.5)
content_placeholder.height = Inches(5)

add_footer(slide, slide_number=2)


# Slide 3: Exploratory Data Analysis (dengan Plot Distribusi Treatment)
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "2. Exploratory Data Analysis", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

# Tambahkan plot Distribusi Treatment
left = Inches(1.5)
top = Inches(2.5)
width = Inches(7)
height = Inches(3.5)
try:
    pic = slide.shapes.add_picture(treatment_plot_path, left, top, width, height)
except FileNotFoundError:
    print(f"Warning: Plot '{treatment_plot_path}' tidak ditemukan. Silakan pastikan file ada di '{plots_dir}'.")
    error_textbox = slide.shapes.add_textbox(left, top + Inches(1), width, Inches(1))
    set_text_properties(error_textbox.text_frame, "Plot Distribusi Treatment tidak dapat dimuat.", Pt(18), color=RGBColor(255, 0, 0), align_center=True)


# Tambahkan caption
caption_textbox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
set_text_properties(caption_textbox.text_frame, "Gambar 1: Distribusi responden berdasarkan status perawatan mental. \nMenunjukkan sebaran yang relatif seimbang (sekitar 51% vs 49%).", Pt(14), color=FOOTER_COLOR, align_center=True)
caption_textbox.top = top + height + Inches(0.2)
caption_textbox.left = left
caption_textbox.width = width
caption_textbox.height = Inches(1) # Beri ruang lebih untuk caption multi-baris

add_footer(slide, slide_number=3)


# Slide 4: Modeling & Evaluasi (dengan Confusion Matrix)
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "3. Modeling & Evaluasi", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

# Tambahkan teks evaluasi
eval_text = (
    f"- Model: Random Forest Classifier\n"
    f"- Skor ROC AUC: {roc_score_value:.2f} → model bekerja sangat baik\n"
    f"- Laporan klasifikasi menunjukkan precision & recall seimbang"
)
eval_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.5))
set_text_properties(eval_textbox.text_frame, eval_text, Pt(20), color=TEXT_COLOR, font_name=BODY_FONT, is_bullet=True)

# Tambahkan plot Confusion Matrix
left = Inches(3) # Posisikan lebih ke tengah
top = Inches(3.5)
width = Inches(4)
height = Inches(3.5)
try:
    pic = slide.shapes.add_picture(confusion_matrix_path, left, top, width, height)
except FileNotFoundError:
    print(f"Warning: Plot '{confusion_matrix_path}' tidak ditemukan. Silakan pastikan file ada di '{plots_dir}'.")
    error_textbox = slide.shapes.add_textbox(left, top + Inches(1), width, height)
    set_text_properties(error_textbox.text_frame, "Plot Confusion Matrix tidak dapat dimuat.", Pt(18), color=RGBColor(255, 0, 0), align_center=True)


# Tambahkan caption
caption_textbox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
set_text_properties(caption_textbox.text_frame, "Gambar 2: Confusion Matrix model Random Forest. \nMenunjukkan akurasi prediksi kelas positif dan negatif.", Pt(14), color=FOOTER_COLOR, align_center=True)
caption_textbox.top = top + height + Inches(0.2)
caption_textbox.left = left
caption_textbox.width = width
caption_textbox.height = Inches(1)

add_footer(slide, slide_number=4)


# Slide 5: Feature Importance (dengan Bar Plot)
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "4. Feature Importance", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

# Konten teks Feature Importance
tf = slide.placeholders[1].text_frame
tf.clear()

p = tf.add_paragraph()
run = p.add_run()
run.text = "Fitur paling berpengaruh:"
run.font.size = Pt(22)
run.font.name = BODY_FONT
run.font.color.rgb = TEXT_COLOR
p.space_after = Pt(6)

for feature in top_features_list_value:
    p_feature = tf.add_paragraph()
    p_feature.level = 0
    
    display_name_map = {
        "work_interfere": "work_interfere (interferensi kerja)",
        "family_history": "family_history (riwayat keluarga)",
        "anonymity": "anonymity (tingkat anonimitas di tempat kerja)"
    }
    
    display_text = display_name_map.get(feature, feature)
    
    bold_part = feature
    rest_part = display_text.replace(bold_part, "").strip()
    
    run_bold = p_feature.add_run()
    run_bold.text = bold_part
    run_bold.font.bold = True
    run_bold.font.size = Pt(22)
    run_bold.font.name = BODY_FONT
    run_bold.font.color.rgb = TEXT_COLOR

    if rest_part:
        run_rest = p_feature.add_run()
        run_rest.text = " " + rest_part
        run_rest.font.bold = False
        run_rest.font.size = Pt(22)
        run_rest.font.name = BODY_FONT
        run_rest.font.color.rgb = TEXT_COLOR
    p_feature.space_after = Pt(4)

p_desc = tf.add_paragraph()
run_desc = p_desc.add_run()
run_desc.text = "\nModel mendeteksi bahwa stres kerja dan dukungan organisasi\nberperan besar terhadap kondisi mental pekerja."
run_desc.font.size = Pt(22)
run_desc.font.name = BODY_FONT
run_desc.font.color.rgb = TEXT_COLOR
p_desc.space_before = Pt(12)

slide.placeholders[1].top = Inches(1.8)
slide.placeholders[1].left = Inches(0.5)
slide.placeholders[1].width = Inches(4.5)
slide.placeholders[1].height = Inches(5)


# Tambahkan plot Feature Importance
left = Inches(5)
top = Inches(1.8)
width = Inches(4.5)
height = Inches(4.5)
try:
    pic = slide.shapes.add_picture(feature_importance_path, left, top, width, height)
except FileNotFoundError:
    print(f"Warning: Plot '{feature_importance_path}' tidak ditemukan. Silakan pastikan file ada di '{plots_dir}'.")
    error_textbox = slide.shapes.add_textbox(left, top + Inches(1), width, height)
    set_text_properties(error_textbox.text_frame, "Plot Feature Importance tidak dapat dimuat.", Pt(18), color=RGBColor(255, 0, 0), align_center=True)


# Tambahkan caption
caption_textbox = slide.shapes.add_textbox(left, top + height + Inches(0.2), width, Inches(0.5))
set_text_properties(caption_textbox.text_frame, "Gambar 3: Bar plot Importansi Fitur. \nMenyoroti fitur 'work_interfere' sebagai yang paling berpengaruh.", Pt(14), color=FOOTER_COLOR, align_center=True)
caption_textbox.top = top + height + Inches(0.2)
caption_textbox.left = left
caption_textbox.width = width
caption_textbox.height = Inches(1)

add_footer(slide, slide_number=5)


# Slide 6: Executive Summary by IBM Granite
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "5. Executive Summary by IBM Granite", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

content_placeholder = slide.placeholders[1]
set_text_properties(content_placeholder.text_frame, final_executive_summary, Pt(24), color=TEXT_COLOR, font_name=BODY_FONT)
content_placeholder.top = Inches(1.8)
content_placeholder.left = Inches(1)
content_placeholder.width = Inches(8)
content_placeholder.height = Inches(5.5)

add_footer(slide, slide_number=6)


# Slide 7: Deployment & Use Case
slide = prs.slides.add_slide(content_slide_layout)
add_simple_background(slide)

set_text_properties(slide.shapes.title.text_frame, "6. Deployment & Use Case", Pt(36), bold=True, color=ACCENT_COLOR, font_name=TITLE_FONT)
slide.shapes.title.top = Inches(0.5)
slide.shapes.title.left = Inches(0.5)

content_placeholder = slide.placeholders[1]
content_text = (
    "- Model dapat diintegrasikan ke sistem HR untuk assessment awal\n"
    "- Bisa dikembangkan untuk chatbot psikologi perusahaan\n"
    "- Visualisasi dan laporan dapat dipresentasikan ke manajemen\n"
    "- Seluruh pipeline berbasis open-source dan cloud-ready"
)
set_text_properties(content_placeholder.text_frame, content_text, Pt(24), color=TEXT_COLOR, font_name=BODY_FONT, is_bullet=True)
content_placeholder.top = Inches(1.8)
content_placeholder.left = Inches(1)
content_placeholder.width = Inches(8.5)
content_placeholder.height = Inches(5)

add_footer(slide, slide_number=7)

# Slide 8: Thank You Slide
slide = prs.slides.add_slide(blank_slide_layout)
add_simple_background(slide, color=ACCENT_COLOR)

txBox = slide.shapes.add_textbox(Inches(0), Inches(0), prs.slide_width, prs.slide_height)
tf = txBox.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

p = tf.paragraphs[0]
run = p.add_run()
run.text = "Terima Kasih!"
font = run.font
font.size = Pt(60)
font.bold = True
font.name = TITLE_FONT
font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

txBox_contact = slide.shapes.add_textbox(Inches(0), Inches(7.5), prs.slide_width, Inches(1))
tf_contact = txBox_contact.text_frame
tf_contact.vertical_anchor = MSO_ANCHOR.MIDDLE

p_contact = tf_contact.paragraphs[0]
run_contact = p_contact.add_run()
run_contact.text = "Untuk pertanyaan lebih lanjut, silakan hubungi tim kami."
font_contact = run_contact.font
font_contact.size = Pt(18)
font_contact.name = BODY_FONT
font_contact.color.rgb = RGBColor(255, 255, 255)
p_contact.alignment = PP_ALIGN.CENTER

add_footer(slide, slide_number=8)


# Simpan presentasi
output_path = "Mental_Health_Capstone_Presentation_Final.pptx"
prs.save(output_path)
print(f"\n✅ Presentasi berhasil disimpan di: {output_path}")
print("--- Pembuatan PPT Selesai ---")

